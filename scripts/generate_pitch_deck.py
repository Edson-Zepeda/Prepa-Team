from __future__ import annotations

import os
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))
sys.path.insert(0, str(ROOT / "src"))

OUT_DIR = ROOT / "presentacion"
ASSET_DIR = OUT_DIR / "assets"
CAPTURE_DIR = ASSET_DIR / "captures"
HTML_DIR = CAPTURE_DIR / "html"
OUTPUT_PPTX = OUT_DIR / "Prepa-Team_Pitch_Final.pptx"

CHROME_CANDIDATES = [
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
]

BG = RGBColor(255, 255, 255)
INK = RGBColor(30, 29, 28)
MUTED = RGBColor(92, 91, 88)
ACCENT = RGBColor(56, 189, 248)
ACCENT_SOFT = RGBColor(224, 246, 254)
LINE = RGBColor(224, 224, 224)


SLIDES = [
    {
        "title": "Sistema de Control Escolar con Intervención Mínima",
        "subtitle": "Prepa-Team · Samsung Innovation Campus 2025-2026",
        "bullets": [
            "Detecta riesgo académico con datos escolares.",
            "Prioriza acciones concretas por estudiante.",
            "Convierte seguimiento escolar en planes accionables.",
        ],
        "tagline": "De predecir rendimiento a activar soluciones.",
        "speech": "Prepa-Team no es solo un modelo predictivo. Es una plataforma de control escolar que ayuda a directivos y estudiantes a identificar riesgo académico y decidir cuál es la intervención mínima más efectiva para mejorar.",
    },
    {
        "title": "Las escuelas detectan tarde el riesgo académico",
        "subtitle": "Cuando el problema ya es visible, muchas veces ya es tarde.",
        "bullets": [
            "Las calificaciones muestran el resultado, no siempre la causa.",
            "El seguimiento académico suele ser manual y reactivo.",
            "No todos los estudiantes necesitan la misma intervención.",
            "La escuela necesita priorizar tutorías y seguimiento.",
        ],
        "tagline": "El riesgo académico necesita acción temprana, no solo reportes.",
        "speech": "El problema no es solo saber quién va mal. El verdadero reto es saber qué hacer primero, con quién hacerlo y qué cambio mínimo puede mover al estudiante a una zona aceptable.",
    },
    {
        "title": "Los datos escolares ya existen",
        "subtitle": "Falta convertirlos en decisiones.",
        "bullets": [
            "Calificaciones, ausencias y hábitos anticipan señales de riesgo.",
            "La escuela puede intervenir antes si ordena esas señales.",
            "La inteligencia artificial ayuda a priorizar acciones.",
            "El valor está en recomendar acciones realistas.",
        ],
        "tagline": "Menos intuición aislada, más intervención basada en evidencia.",
        "speech": "Cada escuela ya genera datos. La oportunidad es usar esos datos para que el directivo no solo vea promedios, sino rutas de acción personalizadas.",
    },
    {
        "title": "Prepa-Team",
        "subtitle": "Sistema de control escolar con recomendaciones personalizadas.",
        "bullets": [
            "Portal para directivos.",
            "Portal para estudiantes.",
            "Carga de calificaciones y reportes.",
            "Planes de intervención mínima.",
        ],
        "tagline": "Una plataforma escolar que no solo informa: guía.",
        "speech": "Prepa-Team une el control escolar tradicional con un motor de intervención. El administrador carga y consulta información; el estudiante ve su desempeño y recibe planes concretos para mejorar.",
        "image": "login.png",
    },
    {
        "title": "No damos consejos genéricos",
        "subtitle": "Calculamos el cambio mínimo necesario.",
        "bullets": [
            "No dice solo: estudia más.",
            "Responde qué cambio exacto ayuda a salir de riesgo.",
            "Compara planes por esfuerzo e impacto.",
            "Prioriza acciones alcanzables.",
        ],
        "tagline": "No más recomendaciones vagas. Planes concretos.",
        "speech": "La diferencia clave es que el sistema no recomienda todo al mismo tiempo. Busca el menor cambio útil para mover al estudiante al siguiente nivel de riesgo.",
        "diagram": "before_after",
    },
    {
        "title": "Dos perfiles, una misma meta",
        "subtitle": "Decisiones escolares más claras.",
        "bullets": [
            "Directivo: gestiona alumnos, materias y calificaciones.",
            "Estudiante: consulta desempeño y plan de mejora.",
            "Escuela: identifica prioridades.",
            "Tutoría: enfoca recursos donde más impactan.",
        ],
        "tagline": "Información útil para quien decide y para quien actúa.",
        "speech": "El sistema está pensado para dos usuarios principales: quien administra la información escolar y quien necesita entender qué puede hacer para mejorar.",
        "diagram": "roles",
    },
    {
        "title": "De la calificación a la intervención",
        "subtitle": "Un flujo simple para uso escolar real.",
        "steps": [
            "Registrar datos",
            "Completar perfil",
            "Calcular riesgo",
            "Generar planes",
            "Priorizar seguimiento",
        ],
        "tagline": "Un proceso claro, repetible y fácil de explicar.",
        "speech": "El flujo está diseñado para una demo realista: cargar datos, consultar desempeño, identificar riesgo y mostrar acciones recomendadas.",
        "diagram": "flow",
    },
    {
        "title": "Panel Directivo",
        "subtitle": "Vista rápida del estado académico.",
        "bullets": [
            "Total de alumnos y promedio general.",
            "Perfiles incompletos y alumnos en riesgo.",
            "Accesos a alumnos, materias y calificaciones.",
            "Seguimiento por alumno.",
        ],
        "tagline": "Visibilidad institucional en una sola pantalla.",
        "speech": "Esta pantalla permite que el directivo tenga una visión general. No reemplaza al tutor, pero ayuda a saber dónde enfocar atención primero.",
        "image": "admin_dashboard.png",
    },
    {
        "title": "Perfil Del Alumno",
        "subtitle": "Datos académicos para personalizar recomendaciones.",
        "bullets": [
            "Edad, horas de estudio y ausencias.",
            "Tutoría, apoyo familiar y actividades.",
            "Variables sensibles fuera de recomendaciones.",
            "Lenguaje claro para estudiantes y directivos.",
        ],
        "tagline": "Personalización sin invadir.",
        "speech": "El sistema usa variables académicas y de contexto escolar. Evitamos presentar variables sensibles en la experiencia del usuario y nos enfocamos en factores accionables.",
        "image": "student_profile.png",
    },
    {
        "title": "Estado Académico",
        "subtitle": "Promedio, probabilidad y nivel actual.",
        "bullets": [
            "Promedio escolar real.",
            "Promedio estimado en escala mexicana.",
            "Probabilidad de buen rendimiento.",
            "Nivel académico y plan principal.",
        ],
        "tagline": "Métricas claras para tomar acción.",
        "speech": "El alumno no ve términos técnicos como GPA. Todo se expresa en escala mexicana de 0 a 10 para que sea entendible y útil.",
        "image": "student_dashboard.png",
    },
    {
        "title": "Plan Mínimo De Mejora",
        "subtitle": "El menor cambio con impacto esperado.",
        "bullets": [
            "Cambio necesario para subir de nivel.",
            "Acciones con cantidad concreta.",
            "Impacto estimado.",
            "Esfuerzo requerido.",
        ],
        "tagline": "La mejor intervención no siempre es la más grande.",
        "speech": "Esta es la parte más importante del proyecto. El sistema no busca que el estudiante cambie todo; busca el cambio mínimo que puede producir una mejora relevante.",
        "image": "student_dashboard.png",
    },
    {
        "title": "Tres Rutas De Acción",
        "subtitle": "La escuela puede elegir según contexto.",
        "bullets": [
            "Menor esfuerzo.",
            "Plan balanceado.",
            "Mayor impacto.",
            "Cada plan muestra acciones, impacto y esfuerzo.",
        ],
        "tagline": "Recomendaciones flexibles, no recetas únicas.",
        "speech": "No todos los estudiantes tienen las mismas posibilidades. Por eso mostramos tres rutas: una de menor esfuerzo, una balanceada y una de mayor impacto.",
        "diagram": "plans",
    },
    {
        "title": "Base Técnica",
        "subtitle": "Modelos simples, interpretables y validados.",
        "bullets": [
            "Dataset académico estructurado.",
            "Preprocesamiento reproducible.",
            "Regresión lineal para promedio estimado.",
            "Clasificador calibrado para buen rendimiento.",
        ],
        "tagline": "Rigor técnico al servicio de decisiones escolares.",
        "speech": "Elegimos modelos interpretables porque el objetivo no es impresionar con complejidad, sino generar recomendaciones confiables y explicables.",
        "diagram": "technical",
    },
    {
        "title": "Validación",
        "subtitle": "El sistema fue probado antes de llevarlo a la web.",
        "bullets": [
            "Regresión estable en datos no vistos.",
            "Clasificación con alta separación de rendimiento.",
            "API validada en escala mexicana 0-10.",
            "Smoke test: login, predicción, reportes y CSV.",
        ],
        "tagline": "La confianza viene de medir, comparar y explicar.",
        "speech": "Validamos que el modelo funciona y también analizamos sus límites. El sistema completo fue probado de punta a punta.",
        "image": "fig_model_comparison.png",
    },
    {
        "title": "Las Ausencias Pesan Mucho",
        "subtitle": "No es solo una variable: es una señal temprana.",
        "bullets": [
            "Al quitar ausencias, el rendimiento del modelo cae.",
            "Esto confirma su importancia como factor de riesgo.",
            "La plataforma convierte esa señal en intervención.",
            "No se interpreta como causa única.",
        ],
        "tagline": "Detectar señales tempranas permite actuar antes.",
        "speech": "El hallazgo no significa que las ausencias expliquen todo. Significa que son una señal muy fuerte para priorizar seguimiento académico.",
        "image": "fig_ablation_absences.png",
    },
    {
        "title": "Arquitectura Del Sistema",
        "subtitle": "Demo funcional lista para escalar.",
        "bullets": [
            "Backend con FastAPI.",
            "SQLite para la demo.",
            "Frontend HTML, CSS y JavaScript.",
            "Control de acceso por rol.",
        ],
        "tagline": "Una demo real, no solo un notebook.",
        "speech": "La demo ya funciona como sistema integrado: login, panel admin, portal alumno, carga de calificaciones y recomendaciones personalizadas.",
        "diagram": "architecture",
    },
    {
        "title": "Impacto Esperado",
        "subtitle": "Mejor priorización, mejor seguimiento.",
        "bullets": [
            "Reduce detección tardía de riesgo.",
            "Ayuda a orientar tutorías.",
            "Da claridad al estudiante.",
            "Puede integrarse a procesos escolares existentes.",
        ],
        "tagline": "Decisiones escolares más rápidas, claras y accionables.",
        "speech": "Para una escuela, el valor no está en predecir por predecir. Está en priorizar recursos, enfocar tutorías y dar al estudiante una ruta clara para mejorar.",
        "diagram": "impact",
    },
    {
        "title": "Próximo Paso: Piloto Escolar",
        "subtitle": "Validar impacto en un entorno real.",
        "bullets": [
            "Ejecutar piloto con un grupo pequeño.",
            "Comparar seguimiento tradicional contra plataforma.",
            "Ajustar mensajes con retroalimentación docente.",
            "Preparar versión escalable institucional.",
        ],
        "tagline": "De demo funcional a herramienta institucional.",
        "speech": "Prepa-Team ya demuestra viabilidad técnica y valor operativo. El siguiente paso natural es un piloto controlado con una escuela para medir utilidad real en seguimiento académico.",
        "diagram": "roadmap",
    },
]


def ensure_dirs() -> None:
    OUT_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)
    CAPTURE_DIR.mkdir(parents=True, exist_ok=True)
    HTML_DIR.mkdir(parents=True, exist_ok=True)


def find_browser() -> Path:
    for candidate in CHROME_CANDIDATES:
        if candidate.exists():
            return candidate
    raise FileNotFoundError("No se encontró Chrome o Edge para generar capturas.")


def inline_css(html: str) -> str:
    css = (ROOT / "web" / "frontend" / "static" / "app.css").read_text(encoding="utf-8")
    replacements = [
        '<link rel="stylesheet" href="http://testserver/static/app.css" />',
        '<link rel="stylesheet" href="/static/app.css" />',
    ]
    for item in replacements:
        html = html.replace(item, f"<style>{css}</style>")
    return html


def write_capture_html(name: str, html: str) -> Path:
    path = HTML_DIR / f"{name}.html"
    path.write_text(inline_css(html), encoding="utf-8")
    return path


def screenshot_html(browser: Path, html_path: Path, image_path: Path, width: int = 1500, height: int = 1000) -> None:
    uri = html_path.resolve().as_uri()
    subprocess.run(
        [
            str(browser),
            "--headless=new",
            "--disable-gpu",
            "--hide-scrollbars",
            f"--window-size={width},{height}",
            f"--screenshot={image_path}",
            uri,
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def generate_web_captures() -> None:
    browser = find_browser()
    with tempfile.TemporaryDirectory(ignore_cleanup_errors=True) as temp_dir:
        os.environ["SCHOOL_DB_PATH"] = str(Path(temp_dir) / "school.db")
        os.environ["SCHOOL_NAME"] = "Sistema de control escolar"

        from fastapi.testclient import TestClient

        from web.backend.app import app

        with TestClient(app) as client:
            pages = {"login": client.get("/login").text}

            admin_login = client.post(
                "/api/auth/login",
                json={"email": "admin@prepateam.local", "password": "Admin1234!"},
            )
            if admin_login.status_code != 200:
                raise RuntimeError(admin_login.text)
            pages["admin_dashboard"] = client.get("/admin/dashboard").text
            pages["admin_student_detail"] = client.get("/admin/students/1").text
            pages["admin_grades_upload"] = client.get("/admin/grades/upload").text

            client.post("/api/auth/logout")
            student_login = client.post(
                "/api/auth/login",
                json={"email": "edson@prepateam.mx", "password": "prepateam"},
            )
            if student_login.status_code != 200:
                raise RuntimeError(student_login.text)
            pages["student_dashboard"] = client.get("/student/dashboard").text
            pages["student_profile"] = client.get("/student/profile").text
            pages["student_report"] = client.get("/student/report").text

        for name, html in pages.items():
            html_path = write_capture_html(name, html)
            screenshot_html(browser, html_path, CAPTURE_DIR / f"{name}.png")


def copy_static_assets() -> None:
    figures = ROOT / "paper" / "figures"
    for name in ["fig_model_comparison.png", "fig_ablation_absences.png", "logo_sic.png"]:
        source = figures / name
        if source.exists():
            (ASSET_DIR / name).write_bytes(source.read_bytes())


def add_notes(slide, notes: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        return


def add_textbox(slide, x, y, w, h, text, size=24, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets: Iterable[str], x, y, w, h, size=21) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.space_after = Pt(7)
    return None


def add_round_rect(slide, x, y, w, h, fill=ACCENT_SOFT, line=LINE, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    return shape


def add_card(slide, x, y, w, h, title, body, fill=ACCENT_SOFT):
    add_round_rect(slide, x, y, w, h, fill=fill)
    add_textbox(slide, x + 0.22, y + 0.18, w - 0.44, 0.35, title, size=17, bold=True)
    add_textbox(slide, x + 0.22, y + 0.62, w - 0.44, h - 0.78, body, size=15, color=MUTED)


def add_header(slide, number: int, title: str, subtitle: str) -> None:
    add_textbox(slide, 0.55, 0.35, 0.45, 0.32, f"{number:02d}", size=13, bold=True, color=ACCENT)
    add_textbox(slide, 0.98, 0.33, 10.2, 0.55, title, size=30, bold=True)
    add_textbox(slide, 1.0, 0.9, 10.2, 0.45, subtitle, size=15, color=MUTED)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.42), Inches(2.1), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def crop_picture(slide, path: Path, x, y, w, h):
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    if pic.height > Inches(h):
        pic.height = Inches(h)
    if pic.width > Inches(w):
        pic.width = Inches(w)
    pic.left = Inches(x + (w - pic.width / 914400) / 2)
    pic.top = Inches(y + (h - pic.height / 914400) / 2)
    return pic


def add_picture_card(slide, image_name: str, x=6.6, y=1.65, w=6.0, h=4.85) -> None:
    add_round_rect(slide, x - 0.06, y - 0.06, w + 0.12, h + 0.12, fill=RGBColor(248, 250, 252), line=LINE)
    path = CAPTURE_DIR / image_name
    if not path.exists():
        path = ASSET_DIR / image_name
    if path.exists():
        crop_picture(slide, path, x, y, w, h)
    else:
        add_textbox(slide, x + 0.2, y + 0.2, w - 0.4, h - 0.4, "Imagen pendiente", size=20, color=MUTED, align=PP_ALIGN.CENTER)


def add_tagline(slide, text: str) -> None:
    add_round_rect(slide, 0.75, 6.65, 11.85, 0.52, fill=INK, line=INK)
    add_textbox(slide, 1.0, 6.75, 11.35, 0.28, text, size=16, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)


def add_flow(slide, labels: list[str], y=3.0) -> None:
    x = 0.9
    for index, label in enumerate(labels):
        add_card(slide, x, y, 2.05, 1.05, f"{index + 1}", label, fill=ACCENT_SOFT)
        if index < len(labels) - 1:
            add_textbox(slide, x + 2.12, y + 0.34, 0.34, 0.25, "→", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        x += 2.38


def add_diagram(slide, kind: str) -> None:
    if kind == "before_after":
        add_card(slide, 6.7, 2.0, 2.55, 2.45, "Antes", "“Estudia más.”\n“Falta menos.”\nConsejos generales.", fill=RGBColor(245, 245, 245))
        add_card(slide, 9.65, 2.0, 2.55, 2.45, "Prepa-Team", "Reducir ausencias de 8 a 4.\nSumar 2 h de estudio.\nImpacto estimado.", fill=ACCENT_SOFT)
        add_textbox(slide, 9.2, 2.95, 0.35, 0.4, "→", size=28, bold=True, color=ACCENT)
    elif kind == "roles":
        add_card(slide, 6.65, 1.95, 2.55, 2.7, "Directivo", "Administra alumnos.\nCarga calificaciones.\nRevisa reportes.", fill=ACCENT_SOFT)
        add_card(slide, 9.65, 1.95, 2.55, 2.7, "Alumno", "Consulta promedio.\nVe factores.\nSigue un plan.", fill=RGBColor(245, 245, 245))
    elif kind == "flow":
        add_flow(slide, ["Registrar datos", "Completar perfil", "Calcular riesgo", "Generar planes", "Priorizar seguimiento"], y=3.0)
    elif kind == "plans":
        add_card(slide, 6.65, 1.8, 1.85, 3.1, "Mínimo", "Menor esfuerzo para subir de nivel.", fill=ACCENT_SOFT)
        add_card(slide, 8.75, 1.8, 1.85, 3.1, "Balanceado", "Mejor relación impacto-esfuerzo.", fill=RGBColor(245, 245, 245))
        add_card(slide, 10.85, 1.8, 1.85, 3.1, "Impacto", "Máxima mejora estimada.", fill=ACCENT_SOFT)
    elif kind == "technical":
        add_flow(slide, ["Datos", "Preproceso", "Modelo", "Riesgo", "Planes"], y=3.1)
    elif kind == "architecture":
        add_flow(slide, ["Usuario", "Web", "FastAPI", "SQLite", "Modelo"], y=3.0)
        add_textbox(slide, 10.5, 4.55, 1.7, 0.5, "Recomendaciones", size=16, bold=True, color=ACCENT)
    elif kind == "impact":
        add_card(slide, 6.65, 1.8, 2.6, 2.8, "Hoy", "Detección tardía.\nRevisión manual.\nRecomendaciones vagas.", fill=RGBColor(245, 245, 245))
        add_card(slide, 9.65, 1.8, 2.6, 2.8, "Con Prepa-Team", "Prioridad clara.\nPlanes concretos.\nSeguimiento medible.", fill=ACCENT_SOFT)
    elif kind == "roadmap":
        add_flow(slide, ["Demo funcional", "Piloto escolar", "Validación", "Escalamiento"], y=3.05)


def build_deck() -> None:
    prs = Presentation()
    prs.slide_width = 18288000
    prs.slide_height = 10287000
    blank = prs.slide_layouts[6]

    for index, item in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG
        add_header(slide, index, item["title"], item["subtitle"])

        if index == 1:
            logo = ASSET_DIR / "logo_sic.png"
            if logo.exists():
                slide.shapes.add_picture(str(logo), Inches(10.9), Inches(0.38), height=Inches(0.55))
            add_textbox(slide, 1.0, 2.0, 5.2, 0.7, "De datos escolares a planes concretos de mejora", size=28, bold=True)
            add_bullets(slide, item["bullets"], 1.05, 3.05, 5.3, 1.9, size=22)
            add_card(slide, 7.1, 2.25, 4.65, 2.35, "Equipo", "Edson Manuel Zepeda Chávez\nFrancisco Ricardo Moreno Sánchez\nAlan Emir Martínez Espinosa", fill=ACCENT_SOFT)
            add_tagline(slide, item["tagline"])
        else:
            if "steps" in item:
                add_diagram(slide, "flow")
            else:
                add_bullets(slide, item["bullets"], 1.05, 1.85, 5.1, 3.75, size=21)

            if "image" in item:
                add_picture_card(slide, item["image"])
            elif "diagram" in item:
                add_diagram(slide, item["diagram"])

            add_tagline(slide, item["tagline"])
        add_notes(slide, item["speech"])

    prs.save(OUTPUT_PPTX)


def validate_deck() -> None:
    prs = Presentation(str(OUTPUT_PPTX))
    if len(prs.slides) != 18:
        raise AssertionError(f"Se esperaban 18 diapositivas y se generaron {len(prs.slides)}")
    text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    banned = ["Proyecto: Predicción de GPA", "Diagrama de Gantt", "Listo para comenzar implementación"]
    found = [item for item in banned if item in text]
    if found:
        raise AssertionError(f"Texto desactualizado encontrado: {found}")
    if "Sistema de Control Escolar con Intervención Mínima" not in text:
        raise AssertionError("No se encontró el título principal esperado.")


def main() -> None:
    ensure_dirs()
    copy_static_assets()
    generate_web_captures()
    build_deck()
    validate_deck()
    print(f"Presentación generada: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
